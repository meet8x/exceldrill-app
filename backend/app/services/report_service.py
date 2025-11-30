from typing import Dict, Any, List
import io
import pandas as pd
import numpy as np
import matplotlib
matplotlib.use('Agg') # Use non-interactive backend
import matplotlib.pyplot as plt
import seaborn as sns
from docx import Document
from docx.shared import Inches, Pt, RGBColor as DocxRGBColor
from pptx import Presentation
from pptx.util import Inches as PptxInches
from pptx.chart.data import CategoryChartData, XyChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from typing import Callable, Optional

class ReportService:
    # Color Schemes
    COLOR_SCHEMES = {
        'kpmg': {
            'primary': '#00338D',
            'secondary': '#0091DA',
            'tertiary': '#005EB8',
            'quaternary': '#483698',
            'background': '#FFFFFF',
            'text': '#000000',
            'palette': ['#00338D', '#0091DA', '#005EB8', '#483698', '#6D2077', '#00A3A1']
        },
        'viridis': {
            'primary': '#440154',
            'secondary': '#3b528b',
            'tertiary': '#21918c',
            'quaternary': '#5ec962',
            'background': '#FFFFFF',
            'text': '#000000',
            'palette': ['#440154', '#3b528b', '#21918c', '#5ec962', '#fde725']
        },
        'seaborn': {
            'primary': '#4C72B0',
            'secondary': '#DD8452',
            'tertiary': '#55A868',
            'quaternary': '#C44E52',
            'background': '#eaeaf2',
            'text': '#333333',
            'palette': ['#4C72B0', '#DD8452', '#55A868', '#C44E52', '#8172B3', '#937860']
        },
        'pastel': {
            'primary': '#A1C9F4',
            'secondary': '#FFB482',
            'tertiary': '#8DE5A1',
            'quaternary': '#FF9F9B',
            'background': '#FFFFFF',
            'text': '#555555',
            'palette': ['#A1C9F4', '#FFB482', '#8DE5A1', '#FF9F9B', '#D0BBFF', '#DEBB9B']
        }
    }

    def __init__(self, color_scheme: str = 'kpmg'):
        self.scheme_name = color_scheme
        self.colors = self.COLOR_SCHEMES.get(color_scheme, self.COLOR_SCHEMES['kpmg'])
        
        # Set matplotlib style based on scheme
        if color_scheme == 'seaborn':
            sns.set_style("whitegrid")
        else:
            sns.set_style("white")
            
        # Update class constants for backward compatibility
        self.KPMG_BLUE = self._hex_to_rgb(self.colors['primary'])
        self.KPMG_LIGHT_BLUE = self._hex_to_rgb(self.colors['secondary'])
    
    def _hex_to_rgb(self, hex_color):
        """Convert hex string to RGBColor object"""
        hex_color = hex_color.lstrip('#')
        return RGBColor(int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16))

    def _truncate_text(self, text: str, max_chars: int = 500) -> str:
        if len(text) <= max_chars:
            return text
        return text[:max_chars] + "..."
    
    def _add_section_separator_slide(self, prs, section_title: str):
        """Add a professional section separator slide"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
        
        # Add title shape
        left = PptxInches(1)
        top = PptxInches(3)
        width = PptxInches(8)
        height = PptxInches(1.5)
        
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.text = section_title
        
        # Format text
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.bold = True
        p.font.size = PptxInches(0.5)
        p.font.color.rgb = self.KPMG_BLUE
        
        # Add background shape
        left = PptxInches(0.5)
        top = PptxInches(2.8)
        width = PptxInches(9)
        height = PptxInches(2)
        
        shape = slide.shapes.add_shape(1, left, top, width, height)  # Rectangle
        fill = shape.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(240, 248, 255)  # Light blue background
        shape.line.color.rgb = self.KPMG_BLUE
        shape.line.width = PptxInches(0.02)

    def _plot_to_bytes(self, fig) -> io.BytesIO:
        buf = io.BytesIO()
        fig.savefig(buf, format='png', bbox_inches='tight')
        buf.seek(0)
        plt.close(fig)
        return buf

    def _generate_univariate_plots(self, df: pd.DataFrame) -> Dict[str, io.BytesIO]:
        plots = {}
        numeric_cols = df.select_dtypes(include=[np.number]).columns
        for col in numeric_cols:
            fig, ax = plt.subplots(figsize=(6, 4))
            sns.histplot(df[col].dropna(), kde=True, ax=ax, color=self.colors['primary'])
            ax.set_title(f'Distribution of {col}', color=self.colors['text'])
            plots[col] = self._plot_to_bytes(fig)
        return plots

    def _generate_correlation_heatmap(self, df: pd.DataFrame) -> io.BytesIO:
        numeric_df = df.select_dtypes(include=[np.number])
        if numeric_df.empty:
            return None
        
        fig, ax = plt.subplots(figsize=(8, 6))
        corr = numeric_df.corr()
        sns.heatmap(corr, annot=True, cmap='coolwarm', fmt=".2f", ax=ax)
        ax.set_title('Correlation Matrix', color=self.colors['text'])
        return self._plot_to_bytes(fig)

    def _generate_bivariate_plots(self, df: pd.DataFrame) -> List[io.BytesIO]:
        plots = []
        numeric_df = df.select_dtypes(include=[np.number])
        if numeric_df.shape[1] < 2:
            return plots

        corr_matrix = numeric_df.corr().abs()
        # Get all pairs
        pairs = (corr_matrix.where(np.triu(np.ones(corr_matrix.shape), k=1).astype(bool))
                 .stack()
                 .sort_values(ascending=False))
        
        # Filter for significant correlation (> 0.3) and limit to top 20 to avoid huge reports
        significant_pairs = pairs[pairs > 0.3].head(20)
        
        for (col1, col2), val in significant_pairs.items():
            fig, ax = plt.subplots(figsize=(6, 4))
            sns.scatterplot(data=df, x=col1, y=col2, ax=ax, color=self.colors['primary'])
            ax.set_title(f'{col1} vs {col2} (Corr: {val:.2f})', color=self.colors['text'])
            plots.append(self._plot_to_bytes(fig))
            
        return plots

    def _add_histogram_slide(self, prs, df, col):
        # Calculate histogram
        data = df[col].dropna()
        if data.empty: return
        
        counts, bin_edges = np.histogram(data, bins='auto')
        
        chart_data = CategoryChartData()
        categories = [f"{bin_edges[i]:.1f}-{bin_edges[i+1]:.1f}" for i in range(len(bin_edges)-1)]
        chart_data.categories = categories
        chart_data.add_series('Count', counts)

        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = self._truncate_text(f"Distribution Analysis: {col}", 50)
        
        # Set title formatting
        title = slide.shapes.title
        title.text_frame.paragraphs[0].font.color.rgb = self.KPMG_BLUE
        title.text_frame.paragraphs[0].font.bold = True
        
        x, y, cx, cy = PptxInches(1), PptxInches(1.5), PptxInches(8), PptxInches(5)
        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
        ).chart
        chart.has_legend = False

    def _add_categorical_slide(self, prs, df, col):
        data = df[col].dropna()
        if data.empty: return
        
        # Get top 10 categories
        counts = data.value_counts().head(10)
        
        # Bar Chart Slide
        chart_data = CategoryChartData()
        chart_data.categories = [self._truncate_text(str(c), 20) for c in counts.index]
        chart_data.add_series('Count', counts.values)

        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = self._truncate_text(f"Category Breakdown: {col}", 50)
        
        # Set title formatting
        title = slide.shapes.title
        title.text_frame.paragraphs[0].font.color.rgb = self.KPMG_BLUE
        
        x, y, cx, cy = PptxInches(1), PptxInches(1.5), PptxInches(8), PptxInches(5)
        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.BAR_CLUSTERED, x, y, cx, cy, chart_data
        ).chart
        chart.has_legend = False

        # Pie Chart Slide (if fewer than 6 categories)
        if len(counts) <= 6:
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            slide.shapes.title.text = self._truncate_text(f"Composition Analysis: {col}", 50)
            
            chart = slide.shapes.add_chart(
                XL_CHART_TYPE.PIE, x, y, cx, cy, chart_data
            ).chart
            chart.has_legend = True

    def _add_scatter_slide(self, prs, df, col1, col2):
        chart_data = XyChartData()
        series = chart_data.add_series(f'{col1} vs {col2}')
        
        # Limit points for performance/file size
        plot_data = df[[col1, col2]].dropna().head(500)
        if plot_data.empty: return

        for _, row in plot_data.iterrows():
            series.add_data_point(row[col1], row[col2])

        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = self._truncate_text(f"Relationship: {col1} vs {col2}", 50)
        
        x, y, cx, cy = PptxInches(1), PptxInches(1.5), PptxInches(8), PptxInches(5)
        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.XY_SCATTER, x, y, cx, cy, chart_data
        ).chart

    def _add_box_plot_slide(self, prs, df, cat_col, num_col):
        """Add a box plot slide for categorical vs numeric analysis"""
        # Group data by category
        grouped_data = df.groupby(cat_col)[num_col].apply(list).to_dict()
        
        # Limit to top 8 categories by count
        top_cats = df[cat_col].value_counts().head(8).index
        filtered_data = {k: v for k, v in grouped_data.items() if k in top_cats}
        
        if not filtered_data:
            return
        
        # Create matplotlib box plot
        fig, ax = plt.subplots(figsize=(8, 5))
        data_to_plot = [filtered_data[cat] for cat in top_cats if cat in filtered_data]
        labels = [self._truncate_text(str(cat), 15) for cat in top_cats if cat in filtered_data]
        
        bp = ax.boxplot(data_to_plot, labels=labels, patch_artist=True)
        
        # Style box plot with selected colors
        for patch in bp['boxes']:
            patch.set_facecolor(self.colors['secondary'])
            patch.set_alpha(0.7)
        
        for whisker in bp['whiskers']:
            whisker.set(color=self.colors['primary'], linewidth=1.5)
        
        for cap in bp['caps']:
            cap.set(color=self.colors['primary'], linewidth=1.5)
        
        for median in bp['medians']:
            median.set(color=self.colors['primary'], linewidth=2)
        
        ax.set_xlabel(cat_col, fontsize=10)
        ax.set_ylabel(num_col, fontsize=10)
        ax.set_title(f'{num_col} by {cat_col}', fontsize=12, fontweight='bold', color=self.colors['text'])
        plt.xticks(rotation=45, ha='right')
        plt.tight_layout()
        
        # Convert to bytes and add to slide
        plot_bytes = self._plot_to_bytes(fig)
        
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = self._truncate_text(f"Comparison: {num_col} by {cat_col}", 50)
        
        # Set title formatting
        title = slide.shapes.title
        title.text_frame.paragraphs[0].font.color.rgb = self.KPMG_BLUE
        title.text_frame.paragraphs[0].font.bold = True
        
        # Add image
        left = PptxInches(0.5)
        top = PptxInches(1.5)
        slide.shapes.add_picture(plot_bytes, left, top, height=PptxInches(5))

    def _detect_variable_type(self, df: pd.DataFrame, col: str) -> str:
        """
        Detects the variable type: 'categorical', 'numerical_discrete', 'numerical_continuous', 'time_series'.
        """
        # 1. Check for Time Series
        if pd.api.types.is_datetime64_any_dtype(df[col]):
            return 'time_series'
        
        # Try parsing as datetime if object
        if df[col].dtype == 'object':
            try:
                pd.to_datetime(df[col], errors='raise')
                pass 
            except:
                pass

        # 2. Check for Numerical
        if pd.api.types.is_numeric_dtype(df[col]):
            # Discrete vs Continuous Heuristic
            # If integers and low cardinality (< 20), treat as discrete
            if pd.api.types.is_integer_dtype(df[col]) and df[col].nunique() < 20:
                return 'numerical_discrete'
            return 'numerical_continuous'

        # 3. Default to Categorical
        return 'categorical'

    def _add_time_series_slide(self, prs, df, col):
        # Resample or aggregate if too many points
        data = df[col].dropna()
        if data.empty: return
        
        try:
            dates = pd.to_datetime(data)
            counts = dates.dt.to_period('M').value_counts().sort_index()
            
            chart_data = CategoryChartData()
            chart_data.categories = [str(d) for d in counts.index]
            chart_data.add_series('Frequency', counts.values)

            slide = prs.slides.add_slide(prs.slide_layouts[5])
            slide.shapes.title.text = self._truncate_text(f"Trend Over Time: {col}", 50)
            
            x, y, cx, cy = PptxInches(1), PptxInches(1.5), PptxInches(8), PptxInches(5)
            chart = slide.shapes.add_chart(
                XL_CHART_TYPE.LINE, x, y, cx, cy, chart_data
            ).chart
            chart.has_legend = False
            self._apply_chart_style(chart)
        except:
            pass

    def _add_frequency_table_slide(self, prs, df, col):
        data = df[col].dropna()
        if data.empty: return
        
        counts = data.value_counts()
        total = len(data)
        percents = (counts / total * 100).round(2) # Rounded to 2 decimals
        
        # Create summary dataframe
        summary = pd.DataFrame({'Count': counts, 'Percentage': percents})
        
        if len(summary) > 12:
            top_10 = summary.head(10)
            others_count = summary.iloc[10:]['Count'].sum()
            others_pct = summary.iloc[10:]['Percentage'].sum()
            # Add 'Others' row using pd.concat
            others_row = pd.DataFrame({'Count': [others_count], 'Percentage': [others_pct]}, index=['Others'])
            summary_display = pd.concat([top_10, others_row])
        else:
            summary_display = summary

        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = self._truncate_text(f"Detailed Breakdown: {col}", 50)
        
        # Add table
        rows, cols = summary_display.shape[0] + 1, 3
        left, top, width, height = PptxInches(1), PptxInches(1.5), PptxInches(8), PptxInches(0.8)
        
        table = slide.shapes.add_table(rows, cols, left, top, width, height).table
        
        # Formatting
        table.columns[0].width = PptxInches(4)
        table.columns[1].width = PptxInches(2)
        table.columns[2].width = PptxInches(2)
        
        # Headers
        table.cell(0, 0).text = "Category"
        table.cell(0, 1).text = "Count"
        table.cell(0, 2).text = "Percentage (%)"
        
        # Data
        for i, (idx, row) in enumerate(summary_display.iterrows()):
            table.cell(i+1, 0).text = self._truncate_text(str(idx), 30)
            table.cell(i+1, 1).text = str(int(row['Count']))
            table.cell(i+1, 2).text = f"{row['Percentage']:.2f}%"
            
            # Font styling
            for j in range(3):
                cell = table.cell(i+1, j)
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.name = 'Calibri'
                    paragraph.font.size = PptxInches(0.14)

    def _apply_chart_style(self, chart):
        # Use selected color palette
        colors = [self._hex_to_rgb(c) for c in self.colors['palette']]
        
        try:
            chart.font.name = 'Calibri'
            chart.font.size = PptxInches(0.12)
            
            # Apply colors to series points if it's a single series (like Bar/Column)
            # Or to series themselves if multiple
            for i, series in enumerate(chart.series):
                fill = series.format.fill
                fill.solid()
                # If it's a pie chart, we need to color points individually
                if chart.chart_type == XL_CHART_TYPE.PIE:
                    for j, point in enumerate(series.points):
                        point.format.fill.solid()
                        point.format.fill.fore_color.rgb = colors[j % len(colors)]
                else:
                    fill.fore_color.rgb = colors[i % len(colors)]
        except:
            pass

    def _is_identifier(self, df, col):
        # 1. Check name for common ID keywords
        lower_col = col.lower()
        if any(x in lower_col for x in ['id', 'email', 'phone', 'mobile', 'uuid', 'guid', 'code', 'token']):
            return True
        
        # 2. Check cardinality for object types (high unique ratio)
        if df[col].dtype == 'object':
            unique_ratio = df[col].nunique() / len(df)
            if unique_ratio > 0.9 and len(df) > 20: 
                return True
        return False

    def _get_dataset_summary(self, df: pd.DataFrame) -> List[str]:
        total_records = len(df)
        total_cells = df.size
        total_missing = df.isnull().sum().sum()
        missing_percent = (total_missing / total_cells) * 100
        
        return [
            f"Total Records: {total_records}",
            f"Total Columns: {len(df.columns)}",
            f"Total Missing Values: {total_missing}",
            f"Overall Missing Percentage: {missing_percent:.2f}%"
        ]

    def _add_identifier_summary_slide(self, prs, df, col):
        data = df[col]
        total = len(data)
        unique = data.nunique()
        missing = data.isnull().sum()
        
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = self._truncate_text(f"Data Summary: {col}", 50)
        
        # Create a simple table
        left, top, width, height = PptxInches(2), PptxInches(2), PptxInches(6), PptxInches(1.5)
        table = slide.shapes.add_table(2, 3, left, top, width, height).table
        
        # Headers
        table.cell(0, 0).text = "Total Records"
        table.cell(0, 1).text = "Unique Values"
        table.cell(0, 2).text = "Missing Values"
        
        # Values
        table.cell(1, 0).text = str(total)
        table.cell(1, 1).text = str(unique)
        table.cell(1, 2).text = str(missing)
        
        # Style
        for i in range(2):
            for j in range(3):
                cell = table.cell(i, j)
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.name = 'Calibri'
                    paragraph.font.size = PptxInches(0.14)

    def generate_word_report(self, df: pd.DataFrame, stats: Dict[str, Any], insights: List[str], progress_callback: Optional[Callable[[int], None]] = None) -> io.BytesIO:
        # Set styles
        doc = Document()
        style = doc.styles['Normal']
        font = style.font
        font.name = 'Calibri'
        font.size = Pt(11)
        font.color.rgb = DocxRGBColor(0, 0, 0) # Default black
        
        from backend.app.services.ai_service import AIService
        ai_service = AIService()

        doc.add_heading('Comprehensive Data Analysis Report', 0)
        
        # 1. Executive Summary
        doc.add_heading('1. Executive Summary', level=1)
        
        # Add Dataset Stats
        doc.add_heading('Dataset Overview', level=2)
        dataset_stats = self._get_dataset_summary(df)
        for stat in dataset_stats:
            doc.add_paragraph(stat, style='List Bullet')
            
        doc.add_heading('Key Insights', level=2)
        for insight in insights:
            doc.add_paragraph(insight, style='List Bullet')
            
        # 2. Statistical Summary
        doc.add_heading('2. Statistical Summary', level=1)
        
        # Data Quality Summary
        doc.add_heading('Data Quality Overview', level=2)
        quality_table = doc.add_table(rows=1, cols=4)
        quality_table.style = 'Table Grid'
        hdr = quality_table.rows[0].cells
        hdr[0].text = 'Column'
        hdr[1].text = 'Missing %'
        hdr[2].text = 'Unique Values'
        hdr[3].text = 'Data Type'
        
        for col in df.columns:
            row = quality_table.add_row().cells
            row[0].text = str(col)
            missing_pct = (df[col].isnull().sum() / len(df) * 100)
            row[1].text = f"{missing_pct:.1f}%"
            row[2].text = str(df[col].nunique())
            row[3].text = str(df[col].dtype)
        
        doc.add_paragraph("")  # Spacer
        
        # Numeric Statistics
        doc.add_heading('Numeric Variables Summary', level=2)
        table = doc.add_table(rows=1, cols=4)
        table.style = 'Table Grid'
        hdr = table.rows[0].cells
        hdr[0].text = 'Column'
        hdr[1].text = 'Mean'
        hdr[2].text = 'Min'
        hdr[3].text = 'Max'
        
        for col, metrics in stats['summary'].items():
            row = table.add_row().cells
            row[0].text = str(col)
            row[1].text = str(round(metrics.get('mean', 0), 2))
            row[2].text = str(metrics.get('min', 0))
            row[3].text = str(metrics.get('max', 0))

        # 3. Univariate Analysis (All Variables)
        doc.add_heading('3. Univariate Analysis', level=1)
        doc.add_paragraph("Analysis of all variables (Numeric & Categorical):")
        
        for i, col in enumerate(df.columns):
            if progress_callback:
                 # Map loop progress to 20-80% range
                 progress = 20 + int((i / len(df.columns)) * 60)
                 progress_callback(progress)

            doc.add_heading(f'Variable: {col}', level=2)

            if self._is_identifier(df, col):
                doc.add_paragraph(f"This variable is identified as an identifier (High cardinality or ID-like name).")
                doc.add_paragraph(f"Unique Values: {df[col].nunique()} | Missing Values: {df[col].isnull().sum()}")
                continue
            
            # Generate AI Insights
            if pd.api.types.is_numeric_dtype(df[col]):
                col_stats = stats['summary'].get(col, {})
                col_stats['count'] = len(df[col].dropna()) # Ensure count is present
                ai_insights = ai_service.generate_variable_insights(col, col_stats, 'numeric')
            else:
                col_stats = df[col].describe().to_dict()
                ai_insights = ai_service.generate_variable_insights(col, col_stats, 'categorical')

            # Add Insights to Word
            doc.add_paragraph(self._truncate_text(ai_insights['short'], 300), style='Intense Quote')
            doc.add_paragraph(self._truncate_text(ai_insights['long'], 1000))

            if pd.api.types.is_numeric_dtype(df[col]):
                # Numeric: Histogram and Box Plot
                fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(10, 4))
                
                # Histogram with KDE
                sns.histplot(df[col].dropna(), kde=True, ax=ax1, color=self.colors['primary'])
                ax1.set_title(f'Distribution of {col}', color=self.colors['text'])
                
                # Box plot for outlier detection
                bp = ax2.boxplot([df[col].dropna()], patch_artist=True)
                bp['boxes'][0].set_facecolor(self.colors['secondary'])
                bp['boxes'][0].set_alpha(0.7)
                bp['medians'][0].set(color=self.colors['primary'], linewidth=2)
                ax2.set_title(f'Outlier Detection: {col}', color=self.colors['text'])
                ax2.set_ylabel(col)
                
                plt.tight_layout()
                doc.add_picture(self._plot_to_bytes(fig), width=Inches(6))
            else:
                # Categorical: Frequency Table
                doc.add_paragraph("Frequency Table:")
                counts = df[col].value_counts()
                percents = (counts / len(df[col].dropna()) * 100).round(1)
                summary = pd.DataFrame({'Count': counts, 'Percentage': percents}).head(10)
                
                ftable = doc.add_table(rows=1, cols=3)
                ftable.style = 'Table Grid'
                hdr = ftable.rows[0].cells
                hdr[0].text = 'Category'
                hdr[1].text = 'Count'
                hdr[2].text = 'Percentage'
                
                for idx, row in summary.iterrows():
                    r = ftable.add_row().cells
                    r[0].text = str(idx)
                    r[1].text = str(int(row['Count']))
                    r[2].text = f"{row['Percentage']}%"
                
                doc.add_paragraph("") # Spacer

                # Categorical: Bar Chart
                fig, ax = plt.subplots(figsize=(6, 4))
                top_cats = df[col].value_counts().head(10)
                sns.barplot(x=top_cats.values, y=[str(x) for x in top_cats.index], ax=ax, palette='viridis')
                ax.set_title(f'Top Categories in {col}', color=self.colors['text'])
                doc.add_picture(self._plot_to_bytes(fig), width=Inches(5))
                
                # Categorical: Pie Chart (if few categories)
                if len(top_cats) <= 6:
                    fig, ax = plt.subplots(figsize=(5, 5))
                    ax.pie(top_cats.values, labels=[str(x) for x in top_cats.index], autopct='%1.1f%%', colors=self.colors['palette'])
                    ax.set_title(f'Proportion of {col}', color=self.colors['text'])
                    doc.add_picture(self._plot_to_bytes(fig), width=Inches(4))

        # 4. Multivariate Analysis (Heatmap)
        if progress_callback: progress_callback(85)
        doc.add_heading('4. Multivariate Analysis', level=1)
        heatmap_bytes = self._generate_correlation_heatmap(df)
        if heatmap_bytes:
            doc.add_picture(heatmap_bytes, width=Inches(6))

        # 5. Key Relationships (Bivariate)
        if progress_callback: progress_callback(90)
        doc.add_heading('5. Key Relationships (Bivariate)', level=1)
        
        # Numeric-Numeric Relationships
        doc.add_heading('Numeric Correlations', level=2)
        bi_plots = self._generate_bivariate_plots(df)
        for plot_bytes in bi_plots:
            doc.add_picture(plot_bytes, width=Inches(5))
        
        # Categorical-Numeric Relationships
        doc.add_heading('Categorical vs Numeric Analysis', level=2)
        categorical_cols = df.select_dtypes(include=['object', 'category']).columns
        numeric_cols = df.select_dtypes(include=[np.number]).columns
        
        # Filter out identifiers
        from backend.app.services.data_processing import DataProcessor
        processor = DataProcessor()
        processor.df = df
        categorical_cols = [col for col in categorical_cols if not processor.is_identifier(col)]
        numeric_cols = [col for col in numeric_cols if not processor.is_identifier(col)]
        
        # Generate box plots for top 5 categorical-numeric pairs
        cat_num_count = 0
        for cat_col in categorical_cols:
            if cat_num_count >= 5:
                break
            for num_col in numeric_cols:
                if cat_num_count >= 5:
                    break
                n_categories = df[cat_col].nunique()
                if 2 <= n_categories <= 8:
                    # Create box plot
                    fig, ax = plt.subplots(figsize=(7, 4))
                    grouped_data = df.groupby(cat_col)[num_col].apply(list).to_dict()
                    top_cats = df[cat_col].value_counts().head(8).index
                    data_to_plot = [grouped_data[cat] for cat in top_cats if cat in grouped_data]
                    labels = [str(cat)[:20] for cat in top_cats if cat in grouped_data]
                    
                    bp = ax.boxplot(data_to_plot, labels=labels, patch_artist=True)
                    for patch in bp['boxes']:
                        patch.set_facecolor(self.colors['secondary'])
                        patch.set_alpha(0.7)
                    for median in bp['medians']:
                        median.set(color=self.colors['primary'], linewidth=2)
                    
                    ax.set_xlabel(cat_col)
                    ax.set_ylabel(num_col)
                    ax.set_title(f'{num_col} by {cat_col}', color=self.colors['text'])
                    plt.xticks(rotation=45, ha='right')
                    plt.tight_layout()
                    
                    doc.add_picture(self._plot_to_bytes(fig), width=Inches(5.5))
                    cat_num_count += 1

        buffer = io.BytesIO()
        doc.save(buffer)
        buffer.seek(0)
        return buffer

    def generate_ppt_report(self, df: pd.DataFrame, stats: Dict[str, Any], insights: List[str], progress_callback: Optional[Callable[[int], None]] = None) -> io.BytesIO:
        prs = Presentation()
        from backend.app.services.ai_service import AIService
        ai_service = AIService()
        
        # Title Slide
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "Data Analysis Report"
        slide.placeholders[1].text = "Automated Comprehensive Analysis"
        
        # Insights & Summary
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Executive Summary"
        tf = slide.placeholders[1].text_frame
        tf.word_wrap = True
        
        # Dataset Stats
        dataset_stats = self._get_dataset_summary(df)
        for stat in dataset_stats:
            p = tf.add_paragraph()
            p.text = stat
            p.level = 0
            p.font.bold = True
            p.font.name = 'Calibri'
            p.font.size = PptxInches(0.16)  # Smaller font
            
        p = tf.add_paragraph()
        p.text = "" # Spacer
        
        # Truncate insights to prevent overflow
        for insight in insights[:3]:  # Limit to 3 for better fit
            p = tf.add_paragraph()
            p.text = self._truncate_text(insight, 150)
            p.level = 0
            p.font.name = 'Calibri'
            p.font.size = PptxInches(0.14)

        # Add Section Separator for Variable Analysis
        self._add_section_separator_slide(prs, "Variable Analysis")

        # Import DataProcessor for identifier detection
        from backend.app.services.data_processing import DataProcessor
        processor = DataProcessor()
        processor.df = df  # Set the dataframe for identifier detection

        # Univariate Analysis (All Variables)
        for i, col in enumerate(df.columns):
            if progress_callback:
                 # Map loop progress to 20-80% range
                 progress = 20 + int((i / len(df.columns)) * 60)
                 progress_callback(progress)
            if processor.is_identifier(col):
                self._add_identifier_summary_slide(prs, df, col)
                continue
            
            var_type = self._detect_variable_type(df, col)
            
            # Generate AI Insights
            if var_type in ['numerical_continuous', 'numerical_discrete']:
                col_stats = stats['summary'].get(col, {})
                col_stats['count'] = len(df[col].dropna())
                ai_insights = ai_service.generate_variable_insights(col, col_stats, 'numeric')
            else:
                col_stats = df[col].describe().to_dict()
                ai_insights = ai_service.generate_variable_insights(col, col_stats, 'categorical')
            
            # Add Insight Slide
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = self._truncate_text(f"Insights: {col}", 50)
            tf = slide.placeholders[1].text_frame
            tf.word_wrap = True
            
            p = tf.add_paragraph()
            p.text = "Summary:"
            p.font.bold = True
            p.font.name = 'Calibri'
            p.font.size = PptxInches(0.16)
            
            p = tf.add_paragraph()
            p.text = self._truncate_text(ai_insights['short'], 120)
            p.level = 1
            p.font.name = 'Calibri'
            p.font.size = PptxInches(0.14)
            
            p = tf.add_paragraph()
            p.text = "Detailed Analysis:"
            p.font.bold = True
            p.font.name = 'Calibri'
            p.font.size = PptxInches(0.16)
            
            p = tf.add_paragraph()
            p.text = self._truncate_text(ai_insights['long'], 250)
            p.level = 1
            p.font.name = 'Calibri'
            p.font.size = PptxInches(0.14)

            # Add Charts based on Type
            if var_type == 'numerical_continuous':
                self._add_histogram_slide(prs, df, col)
                self._apply_chart_style(prs.slides[-1].shapes[-1].chart)
            elif var_type == 'numerical_discrete':
                # For discrete, a bar chart of counts might be better than a histogram if few values
                self._add_categorical_slide(prs, df, col) # Re-use bar chart logic
                self._apply_chart_style(prs.slides[-1].shapes[-1].chart)
            elif var_type == 'time_series':
                self._add_time_series_slide(prs, df, col)
            else: # Categorical
                self._add_frequency_table_slide(prs, df, col)
                self._add_categorical_slide(prs, df, col)
                self._apply_chart_style(prs.slides[-1].shapes[-1].chart)

        # Add Section Separator for Relationship Analysis
        if progress_callback: progress_callback(85)
        self._add_section_separator_slide(prs, "Relationship Analysis")

        # Heatmap
        heatmap_bytes = self._generate_correlation_heatmap(df)
        if heatmap_bytes:
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            slide.shapes.title.text = "Correlation Matrix"
            slide.shapes.add_picture(heatmap_bytes, PptxInches(1), PptxInches(1.5), height=PptxInches(5))

        # Bivariate Plots (Numeric vs Numeric)
        numeric_df = df.select_dtypes(include=[np.number])
        
        # Filter out identifiers from numeric columns
        from backend.app.services.data_processing import DataProcessor
        processor = DataProcessor()
        processor.df = df
        numeric_cols_filtered = [col for col in numeric_df.columns if not processor.is_identifier(col)]
        numeric_df = numeric_df[numeric_cols_filtered]
        
        if numeric_df.shape[1] >= 2:
            corr_matrix = numeric_df.corr().abs()
            pairs = (corr_matrix.where(np.triu(np.ones(corr_matrix.shape), k=1).astype(bool))
                     .stack()
                     .sort_values(ascending=False)
                     .head(20)) # Limit to top 20 significant
            
            for (col1, col2), val in pairs.items():
                self._add_scatter_slide(prs, df, col1, col2)
                self._apply_chart_style(prs.slides[-1].shapes[-1].chart)
        
        # Categorical vs Numeric Analysis
        categorical_cols = df.select_dtypes(include=['object', 'category']).columns
        numeric_cols = df.select_dtypes(include=[np.number]).columns
        
        # Filter out identifiers
        from backend.app.services.data_processing import DataProcessor
        processor = DataProcessor()
        processor.df = df
        categorical_cols = [col for col in categorical_cols if not processor.is_identifier(col)]
        numeric_cols = [col for col in numeric_cols if not processor.is_identifier(col)]
        
        # Generate box plots for categorical vs numeric (limit to top 10 combinations)
        cat_num_pairs = []
        for cat_col in categorical_cols:
            for num_col in numeric_cols:
                # Only include if categorical has reasonable number of categories (2-10)
                n_categories = df[cat_col].nunique()
                if 2 <= n_categories <= 10:
                    cat_num_pairs.append((cat_col, num_col))
        
        # Limit to 10 most interesting pairs
        for cat_col, num_col in cat_num_pairs[:10]:
            self._add_box_plot_slide(prs, df, cat_col, num_col)

        buffer = io.BytesIO()
        prs.save(buffer)
        buffer.seek(0)
        return buffer
