from typing import Dict, Any, List
import io
import pandas as pd
import numpy as np
import matplotlib
matplotlib.use('Agg') # Use non-interactive backend
import matplotlib.pyplot as plt
import seaborn as sns
from docx import Document
from docx.shared import Inches
from pptx import Presentation
from pptx.util import Inches as PptxInches
from pptx.chart.data import CategoryChartData, XyChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.dml.color import RGBColor

class ReportService:
    def _plot_to_bytes(self, fig) -> io.BytesIO:
        buf = io.BytesIO()
        fig.savefig(buf, format='png', bbox_inches='tight')
        buf.seek(0)
        plt.close(fig)
        return buf

    def _generate_univariate_plots(self, df: pd.DataFrame) -> Dict[str, io.BytesIO]:
        plots = {}
        numeric_cols = df.select_dtypes(include=[np.number]).columns
        for col in numeric_cols:
            fig, ax = plt.subplots(figsize=(6, 4))
            sns.histplot(df[col].dropna(), kde=True, ax=ax)
            ax.set_title(f'Distribution of {col}')
            plots[col] = self._plot_to_bytes(fig)
        return plots

    def _generate_correlation_heatmap(self, df: pd.DataFrame) -> io.BytesIO:
        numeric_df = df.select_dtypes(include=[np.number])
        if numeric_df.empty:
            return None
        
        fig, ax = plt.subplots(figsize=(8, 6))
        corr = numeric_df.corr()
        sns.heatmap(corr, annot=True, cmap='coolwarm', fmt=".2f", ax=ax)
        ax.set_title('Correlation Matrix')
        return self._plot_to_bytes(fig)

    def _generate_bivariate_plots(self, df: pd.DataFrame) -> List[io.BytesIO]:
        plots = []
        numeric_df = df.select_dtypes(include=[np.number])
        if numeric_df.shape[1] < 2:
            return plots

        corr_matrix = numeric_df.corr().abs()
        # Get top 5 correlated pairs (excluding self-correlation)
        pairs = (corr_matrix.where(np.triu(np.ones(corr_matrix.shape), k=1).astype(bool))
                 .stack()
                 .sort_values(ascending=False)
                 .head(5))
        
        for (col1, col2), val in pairs.items():
            fig, ax = plt.subplots(figsize=(6, 4))
            sns.scatterplot(data=df, x=col1, y=col2, ax=ax)
            ax.set_title(f'{col1} vs {col2} (Corr: {val:.2f})')
            plots.append(self._plot_to_bytes(fig))
            
        return plots

    def _add_histogram_slide(self, prs, df, col):
        # Calculate histogram
        data = df[col].dropna()
        if data.empty: return
        
        counts, bin_edges = np.histogram(data, bins='auto')
        
        chart_data = CategoryChartData()
        categories = [f"{bin_edges[i]:.1f}-{bin_edges[i+1]:.1f}" for i in range(len(bin_edges)-1)]
        chart_data.categories = categories
        chart_data.add_series('Count', counts)

        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = f"Distribution (Numeric): {col}"
        
        x, y, cx, cy = PptxInches(1), PptxInches(1.5), PptxInches(8), PptxInches(5)
        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
        ).chart
        chart.has_legend = False

    def _add_categorical_slide(self, prs, df, col):
        data = df[col].dropna()
        if data.empty: return
        
        # Get top 10 categories
        counts = data.value_counts().head(10)
        
        # Bar Chart Slide
        chart_data = CategoryChartData()
        chart_data.categories = [str(c) for c in counts.index]
        chart_data.add_series('Count', counts.values)

        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = f"Frequency (Categorical): {col}"
        
        x, y, cx, cy = PptxInches(1), PptxInches(1.5), PptxInches(8), PptxInches(5)
        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.BAR_CLUSTERED, x, y, cx, cy, chart_data
        ).chart
        chart.has_legend = False

        # Pie Chart Slide (if fewer than 6 categories)
        if len(counts) <= 6:
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            slide.shapes.title.text = f"Proportion (Categorical): {col}"
            
            chart = slide.shapes.add_chart(
                XL_CHART_TYPE.PIE, x, y, cx, cy, chart_data
            ).chart
            chart.has_legend = True

    def _add_scatter_slide(self, prs, df, col1, col2):
        chart_data = XyChartData()
        series = chart_data.add_series(f'{col1} vs {col2}')
        
        # Limit points for performance/file size
        plot_data = df[[col1, col2]].dropna().head(500)
        if plot_data.empty: return

        for _, row in plot_data.iterrows():
            series.add_data_point(row[col1], row[col2])

        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = f"{col1} vs {col2}"
        
        x, y, cx, cy = PptxInches(1), PptxInches(1.5), PptxInches(8), PptxInches(5)
        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.XY_SCATTER, x, y, cx, cy, chart_data
        ).chart

    def _detect_variable_type(self, df: pd.DataFrame, col: str) -> str:
        """
        Detects the variable type: 'categorical', 'numerical_discrete', 'numerical_continuous', 'time_series'.
        """
        # 1. Check for Time Series
        if pd.api.types.is_datetime64_any_dtype(df[col]):
            return 'time_series'
        
        # Try parsing as datetime if object
        if df[col].dtype == 'object':
            try:
                pd.to_datetime(df[col], errors='raise')
                # If successful for a sample, assume time series (but be careful of false positives like simple numbers)
                # For safety, only if it really looks like a date. 
                # Let's rely on explicit conversion in data_processing or just skip for now to avoid errors.
                pass 
            except:
                pass

        # 2. Check for Numerical
        if pd.api.types.is_numeric_dtype(df[col]):
            # Discrete vs Continuous Heuristic
            # If integers and low cardinality (< 20), treat as discrete
            if pd.api.types.is_integer_dtype(df[col]) and df[col].nunique() < 20:
                return 'numerical_discrete'
            return 'numerical_continuous'

        # 3. Default to Categorical
        return 'categorical'

    def _add_time_series_slide(self, prs, df, col):
        # Resample or aggregate if too many points
        data = df[col].dropna()
        if data.empty: return
        
        # For a single column, a line chart of values vs index (if index is time) or just values in order?
        # Usually Time Series implies (Time Column, Value Column). 
        # If 'col' IS the time column, we can't plot it alone as a line chart (it's the X axis).
        # If 'col' is a value and there is a time index, we plot value over time.
        # For univariate of a datetime column, maybe a histogram of dates (Frequency over time)?
        
        # Let's assume Univariate Time Series Analysis = Frequency of events over time (Histogram of dates)
        try:
            dates = pd.to_datetime(data)
            counts = dates.dt.to_period('M').value_counts().sort_index()
            
            chart_data = CategoryChartData()
            chart_data.categories = [str(d) for d in counts.index]
            chart_data.add_series('Frequency', counts.values)

            slide = prs.slides.add_slide(prs.slide_layouts[5])
            slide.shapes.title.text = f"Time Trend (Frequency): {col}"
            
            x, y, cx, cy = PptxInches(1), PptxInches(1.5), PptxInches(8), PptxInches(5)
            chart = slide.shapes.add_chart(
                XL_CHART_TYPE.LINE, x, y, cx, cy, chart_data
            ).chart
            chart.has_legend = False
            self._apply_chart_style(chart)
        except:
            pass

    def _add_frequency_table_slide(self, prs, df, col):
        data = df[col].dropna()
        if data.empty: return
        
        counts = data.value_counts()
        total = len(data)
        percents = (counts / total * 100).round(2) # Rounded to 2 decimals
        
        # Create summary dataframe
        summary = pd.DataFrame({'Count': counts, 'Percentage': percents})
        
        # Split logic: If > 10 categories, show Top 10 on one slide, maybe mention others?
        # User asked to split slides if dense. 
        # Strategy: Always show Top 10. If more, create a second slide for "Next 10" or just summarize "Others".
        # For simplicity and clarity in automated reports, Top 10 + Others is best standard.
        
        if len(summary) > 12:
            top_10 = summary.head(10)
            others_count = summary.iloc[10:]['Count'].sum()
            others_pct = summary.iloc[10:]['Percentage'].sum()
            # Add 'Others' row using pd.concat
            others_row = pd.DataFrame({'Count': [others_count], 'Percentage': [others_pct]}, index=['Others'])
            summary_display = pd.concat([top_10, others_row])
        else:
            summary_display = summary

        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = f"Frequency Table: {col}"
        
        # Add table
        rows, cols = summary_display.shape[0] + 1, 3
        left, top, width, height = PptxInches(1), PptxInches(1.5), PptxInches(8), PptxInches(0.8)
        
        table = slide.shapes.add_table(rows, cols, left, top, width, height).table
        
        # Formatting
        table.columns[0].width = PptxInches(4)
        table.columns[1].width = PptxInches(2)
        table.columns[2].width = PptxInches(2)
        
        # Headers
        table.cell(0, 0).text = "Category"
        table.cell(0, 1).text = "Count"
        table.cell(0, 2).text = "Percentage (%)"
        
        # Data
        for i, (idx, row) in enumerate(summary_display.iterrows()):
            table.cell(i+1, 0).text = str(idx)
            table.cell(i+1, 1).text = str(int(row['Count']))
            table.cell(i+1, 2).text = f"{row['Percentage']:.2f}%"
            
            # Font styling
            for j in range(3):
                cell = table.cell(i+1, j)
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.name = 'Calibri'
                    paragraph.font.size = PptxInches(0.14)

    def _apply_chart_style(self, chart):
        # Professional Big-4 Palette (Blended)
        # Blue, Red, Green, Yellow, Grey
        # Hex codes: ['#00338D', '#483698', '#470A68', '#0091DA', '#6D2077', '#005EB8']
        colors = [
            RGBColor(0, 51, 141),    # #00338D
            RGBColor(72, 54, 152),   # #483698
            RGBColor(71, 10, 104),   # #470A68
            RGBColor(0, 145, 218),   # #0091DA
            RGBColor(109, 32, 119),  # #6D2077
            RGBColor(0, 94, 184)     # #005EB8
        ]
        
        try:
            chart.font.name = 'Calibri'
            chart.font.size = PptxInches(0.12)
            
            # Apply colors to series points if it's a single series (like Bar/Column)
            # Or to series themselves if multiple
            for i, series in enumerate(chart.series):
                fill = series.format.fill
                fill.solid()
                # If it's a pie chart, we need to color points individually
                if chart.chart_type == XL_CHART_TYPE.PIE:
                    for j, point in enumerate(series.points):
                        point.format.fill.solid()
                        point.format.fill.fore_color.rgb = colors[j % len(colors)]
                else:
                    fill.fore_color.rgb = colors[i % len(colors)]
        except:
            pass

    def _is_identifier(self, df, col):
        # 1. Check name for common ID keywords
        lower_col = col.lower()
        if any(x in lower_col for x in ['id', 'email', 'phone', 'mobile', 'uuid', 'guid', 'code', 'token']):
            return True
        
        # 2. Check cardinality for object types (high unique ratio)
        if df[col].dtype == 'object':
            unique_ratio = df[col].nunique() / len(df)
            if unique_ratio > 0.9 and len(df) > 20: 
                return True
        return False

    def _get_dataset_summary(self, df: pd.DataFrame) -> List[str]:
        total_records = len(df)
        total_cells = df.size
        total_missing = df.isnull().sum().sum()
        missing_percent = (total_missing / total_cells) * 100
        
        return [
            f"Total Records: {total_records}",
            f"Total Columns: {len(df.columns)}",
            f"Total Missing Values: {total_missing}",
            f"Overall Missing Percentage: {missing_percent:.2f}%"
        ]

    def _add_identifier_summary_slide(self, prs, df, col):
        data = df[col]
        total = len(data)
        unique = data.nunique()
        missing = data.isnull().sum()
        
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = f"Summary: {col} (Identifier)"
        
        # Create a simple table
        left, top, width, height = PptxInches(2), PptxInches(2), PptxInches(6), PptxInches(1.5)
        table = slide.shapes.add_table(2, 3, left, top, width, height).table
        
        # Headers
        table.cell(0, 0).text = "Total Records"
        table.cell(0, 1).text = "Unique Values"
        table.cell(0, 2).text = "Missing Values"
        
        # Values
        table.cell(1, 0).text = str(total)
        table.cell(1, 1).text = str(unique)
        table.cell(1, 2).text = str(missing)
        
        # Style
        for i in range(2):
            for j in range(3):
                cell = table.cell(i, j)
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.name = 'Calibri'
                    paragraph.font.size = PptxInches(0.14)

    def generate_word_report(self, df: pd.DataFrame, stats: Dict[str, Any], insights: List[str]) -> io.BytesIO:
        # Set Professional Style (KPMG-like: Blue/Grey)
        plt.style.use('seaborn-v0_8-whitegrid')
        plt.rcParams['axes.prop_cycle'] = plt.cycler(color=['#00338D', '#483698', '#470A68', '#0091DA', '#6D2077', '#005EB8'])
        
        from backend.app.services.ai_service import AIService
        ai_service = AIService()

        doc = Document()
        doc.add_heading('Comprehensive Data Analysis Report', 0)
        
        # 1. Executive Summary
        doc.add_heading('1. Executive Summary', level=1)
        
        # Add Dataset Stats
        doc.add_heading('Dataset Overview', level=2)
        dataset_stats = self._get_dataset_summary(df)
        for stat in dataset_stats:
            doc.add_paragraph(stat, style='List Bullet')
            
        doc.add_heading('Key Insights', level=2)
        for insight in insights:
            doc.add_paragraph(insight, style='List Bullet')
            
        # 2. Statistical Summary
        doc.add_heading('2. Statistical Summary', level=1)
        table = doc.add_table(rows=1, cols=4)
        table.style = 'Table Grid'
        hdr = table.rows[0].cells
        hdr[0].text = 'Column'
        hdr[1].text = 'Mean'
        hdr[2].text = 'Min'
        hdr[3].text = 'Max'
        
        for col, metrics in stats['summary'].items():
            row = table.add_row().cells
            row[0].text = str(col)
            row[1].text = str(round(metrics.get('mean', 0), 2))
            row[2].text = str(metrics.get('min', 0))
            row[3].text = str(metrics.get('max', 0))

        # 3. Univariate Analysis (All Variables)
        doc.add_heading('3. Univariate Analysis', level=1)
        doc.add_paragraph("Analysis of all variables (Numeric & Categorical):")
        
        for col in df.columns:
            doc.add_heading(f'Variable: {col}', level=2)

            if self._is_identifier(df, col):
                doc.add_paragraph(f"This variable is identified as an identifier (High cardinality or ID-like name).")
                doc.add_paragraph(f"Unique Values: {df[col].nunique()} | Missing Values: {df[col].isnull().sum()}")
                continue
            
            # Generate AI Insights
            if pd.api.types.is_numeric_dtype(df[col]):
                col_stats = stats['summary'].get(col, {})
                col_stats['count'] = len(df[col].dropna()) # Ensure count is present
                ai_insights = ai_service.generate_variable_insights(col, col_stats, 'numeric')
            else:
                col_stats = df[col].describe().to_dict()
                ai_insights = ai_service.generate_variable_insights(col, col_stats, 'categorical')

            # Add Insights to Word
            doc.add_paragraph(ai_insights['short'], style='Intense Quote')
            doc.add_paragraph(ai_insights['long'])

            if pd.api.types.is_numeric_dtype(df[col]):
                # Numeric: Histogram
                fig, ax = plt.subplots(figsize=(6, 4))
                sns.histplot(df[col].dropna(), kde=True, ax=ax, color='#00338D') # KPMG Blue
                ax.set_title(f'Distribution of {col}')
                doc.add_picture(self._plot_to_bytes(fig), width=Inches(5))
            else:
                # Categorical: Frequency Table
                doc.add_paragraph("Frequency Table:")
                counts = df[col].value_counts()
                percents = (counts / len(df[col].dropna()) * 100).round(1)
                summary = pd.DataFrame({'Count': counts, 'Percentage': percents}).head(10)
                
                ftable = doc.add_table(rows=1, cols=3)
                ftable.style = 'Table Grid'
                hdr = ftable.rows[0].cells
                hdr[0].text = 'Category'
                hdr[1].text = 'Count'
                hdr[2].text = 'Percentage'
                
                for idx, row in summary.iterrows():
                    r = ftable.add_row().cells
                    r[0].text = str(idx)
                    r[1].text = str(int(row['Count']))
                    r[2].text = f"{row['Percentage']}%"
                
                doc.add_paragraph("") # Spacer

                # Categorical: Bar Chart
                fig, ax = plt.subplots(figsize=(6, 4))
                top_cats = df[col].value_counts().head(10)
                sns.barplot(x=top_cats.values, y=[str(x) for x in top_cats.index], ax=ax, palette='viridis')
                ax.set_title(f'Top Categories in {col}')
                doc.add_picture(self._plot_to_bytes(fig), width=Inches(5))
                
                # Categorical: Pie Chart (if few categories)
                if len(top_cats) <= 6:
                    fig, ax = plt.subplots(figsize=(5, 5))
                    ax.pie(top_cats.values, labels=[str(x) for x in top_cats.index], autopct='%1.1f%%', colors=['#00338D', '#0091DA', '#6D2077', '#005EB8', '#483698'])
                    ax.set_title(f'Proportion of {col}')
                    doc.add_picture(self._plot_to_bytes(fig), width=Inches(4))

        # 4. Multivariate Analysis (Heatmap)
        doc.add_heading('4. Multivariate Analysis', level=1)
        heatmap_bytes = self._generate_correlation_heatmap(df)
        if heatmap_bytes:
            doc.add_picture(heatmap_bytes, width=Inches(6))

        # 5. Bivariate Analysis (Top Correlations)
        doc.add_heading('5. Key Relationships (Bivariate)', level=1)
        bi_plots = self._generate_bivariate_plots(df)
        for plot_bytes in bi_plots:
            doc.add_picture(plot_bytes, width=Inches(5))

        buffer = io.BytesIO()
        doc.save(buffer)
        buffer.seek(0)
        return buffer

    def generate_ppt_report(self, df: pd.DataFrame, stats: Dict[str, Any], insights: List[str]) -> io.BytesIO:
        prs = Presentation()
        from backend.app.services.ai_service import AIService
        ai_service = AIService()
        
        # Title Slide
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "Data Analysis Report"
        slide.placeholders[1].text = "Automated Comprehensive Analysis"
        
        # Insights & Summary
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Executive Summary"
        tf = slide.placeholders[1].text_frame
        
        # Dataset Stats
        dataset_stats = self._get_dataset_summary(df)
        for stat in dataset_stats:
            p = tf.add_paragraph()
            p.text = stat
            p.level = 0
            p.font.bold = True
            p.font.name = 'Calibri'
            
        p = tf.add_paragraph()
        p.text = "" # Spacer
        
        for insight in insights[:4]: # Limit to fit
            p = tf.add_paragraph()
            p.text = insight
            p.level = 0
            p.font.name = 'Calibri'

        # Univariate Analysis (All Variables)
        for col in df.columns:
            if self._is_identifier(df, col):
                self._add_identifier_summary_slide(prs, df, col)
                continue
            
            var_type = self._detect_variable_type(df, col)
            
            # Generate AI Insights
            if var_type in ['numerical_continuous', 'numerical_discrete']:
                col_stats = stats['summary'].get(col, {})
                col_stats['count'] = len(df[col].dropna())
                ai_insights = ai_service.generate_variable_insights(col, col_stats, 'numeric')
            else:
                col_stats = df[col].describe().to_dict()
                ai_insights = ai_service.generate_variable_insights(col, col_stats, 'categorical')
            
            # Add Insight Slide
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = f"Insights: {col}"
            tf = slide.placeholders[1].text_frame
            
            p = tf.add_paragraph()
            p.text = "Summary:"
            p.font.bold = True
            p.font.name = 'Calibri'
            
            p = tf.add_paragraph()
            p.text = ai_insights['short']
            p.level = 1
            p.font.name = 'Calibri'
            
            p = tf.add_paragraph()
            p.text = "Detailed Analysis:"
            p.font.bold = True
            p.font.name = 'Calibri'
            
            p = tf.add_paragraph()
            p.text = ai_insights['long']
            p.level = 1
            p.font.name = 'Calibri'

            # Add Charts based on Type
            if var_type == 'numerical_continuous':
                self._add_histogram_slide(prs, df, col)
                self._apply_chart_style(prs.slides[-1].shapes[-1].chart)
            elif var_type == 'numerical_discrete':
                # For discrete, a bar chart of counts might be better than a histogram if few values
                self._add_categorical_slide(prs, df, col) # Re-use bar chart logic
                self._apply_chart_style(prs.slides[-1].shapes[-1].chart)
            elif var_type == 'time_series':
                self._add_time_series_slide(prs, df, col)
            else: # Categorical
                self._add_frequency_table_slide(prs, df, col)
                self._add_categorical_slide(prs, df, col)
                self._apply_chart_style(prs.slides[-1].shapes[-1].chart)

        # Heatmap
        heatmap_bytes = self._generate_correlation_heatmap(df)
        if heatmap_bytes:
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            slide.shapes.title.text = "Correlation Matrix"
            slide.shapes.add_picture(heatmap_bytes, PptxInches(1), PptxInches(1.5), height=PptxInches(5))

        # Bivariate Plots
        numeric_df = df.select_dtypes(include=[np.number])
        if numeric_df.shape[1] >= 2:
            corr_matrix = numeric_df.corr().abs()
            pairs = (corr_matrix.where(np.triu(np.ones(corr_matrix.shape), k=1).astype(bool))
                     .stack()
                     .sort_values(ascending=False)
                     .head(5))
            
            for (col1, col2), val in pairs.items():
                self._add_scatter_slide(prs, df, col1, col2)
                self._apply_chart_style(prs.slides[-1].shapes[-1].chart)

        buffer = io.BytesIO()
        prs.save(buffer)
        buffer.seek(0)
        return buffer
