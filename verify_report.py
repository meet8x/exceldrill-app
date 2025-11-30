import requests
import time
import subprocess
import sys
import os
from pptx import Presentation
from pptx.enum.chart import XL_CHART_TYPE

def verify_report():
    # Assume Backend is running
    print("Verifying report generation...", flush=True)
    
    base_url = "http://localhost:8000/api"
    
    try:
        # 1. Upload File
        print("Uploading file...", flush=True)
        try:
            files = {'file': ('dummy.csv', open('dummy.csv', 'rb'), 'text/csv')}
            response = requests.post(f"{base_url}/upload", files=files)
        except requests.exceptions.ConnectionError:
            print("FAIL: Could not connect to backend. Is it running?", flush=True)
            return False

        if response.status_code != 200:
            print(f"Upload failed: {response.text}", flush=True)
            return False
            
        # 2. Generate Report
        print("Generating report...", flush=True)
        response = requests.get(f"{base_url}/report/ppt")
        if response.status_code != 200:
            print(f"Report generation failed: {response.text}", flush=True)
            return False
            
        with open("report.pptx", "wb") as f:
            f.write(response.content)
            
        print("Report saved to report.pptx", flush=True)
        
        # 3. Verify Content
        print("Verifying report content...", flush=True)
        prs = Presentation("report.pptx")
        
        slides_titles = []
        for slide in prs.slides:
            if slide.shapes.title:
                slides_titles.append(slide.shapes.title.text)
        
        print("Slides found:", slides_titles, flush=True)
        
        # Check Executive Summary
        if "Executive Summary" not in slides_titles:
            print("FAIL: Executive Summary missing", flush=True)
            return False
            
        # Check Identifier Summary
        if "Summary: id (Identifier)" not in slides_titles:
            print("FAIL: Identifier Summary for 'id' missing", flush=True)
            return False
            
        # Check Frequency Table
        freq_slide_found = False
        for title in slides_titles:
            if "Frequency Table" in title:
                freq_slide_found = True
                break
        if not freq_slide_found:
            print("FAIL: Frequency Table missing", flush=True)
            return False
            
        print("Verification Successful!", flush=True)
        return True
        
    except Exception as e:
        print(f"An error occurred: {e}", flush=True)
        return False

if __name__ == "__main__":
    if verify_report():
        sys.exit(0)
    else:
        sys.exit(1)
